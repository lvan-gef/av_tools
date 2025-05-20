import argparse
import sys
from pathlib import Path
from collections import namedtuple
import shutil

import pymupdf
from pptx import Presentation
from pptx.util import Inches

RESOLUTION = namedtuple('resolution', ['width', 'height'])


def main(pdf: Path, reso: namedtuple, out: Path) -> None:
    outdir = pdf.parent.joinpath('.converted')
    try:
        outdir.mkdir(parents=True, exist_ok=True)
    except PermissionError:
        print(f'You don\'t have permission to path: "{outdir}"', file=sys.stderr)
        exit(4)
    except Exception as e:
        print(f'A unexpected error: "{e}", while creating path: {outdir}', file=sys.stderr)
        exit(5)

    try:
        pngs = _pdf_to_png(pdf=pdf, reso=reso, outdir=outdir)
        pptx = _png_to_pptx(pngs_list=pngs, pptx=out)

        shutil.move(pptx, out)
        print(f'Powerpoint is saved at location: "{out}"')
    except PermissionError as pe:
        print(pe, file=sys.stderr)
    except ValueError as ve:
        print(ve, file=sys.stderr)
    except FileNotFoundError as fe:
        print(fe, file=sys.stderr)
    except IOError as ioe:
        print(ioe, file=sys.stderr)
    except RuntimeError as re:
        print(re, file=sys.stderr)
    except Exception as e:
        print(e, file=sys.stderr)
    finally:
        if outdir.exists():
            try:
                shutil.rmtree(str(outdir))
            except Exception as e:
                print(f"Warning: Failed to clean up temporary files: {e}", file=sys.stderr)


def _pdf_to_png(pdf: Path, reso: namedtuple, outdir: Path) -> list[Path]:
    outlist_pngs = []

    print('Start parsing pdf')
    try:
        doc = pymupdf.open(pdf)
    except PermissionError:
        raise PermissionError(f'You don\'t have permission to path: "{pdf}"')
    except Exception as e:
        raise RuntimeError(f'A unexpected error: "{e}", while opening pdf: {outdir}')

    try:
        if len(doc) == 0:
            raise ValueError(f'pdf: "{pdf}", is empty')

        for i, page in enumerate(doc, start=1):
            try:
                rect = page.rect
                page_width, page_height = rect.width, rect.height
                if page_width <= 0 or page_height <= 0:
                    raise ValueError(f'page: {i}, have a invalid dimension')

                scale_width = reso.width / page_width
                scale_height = reso.height / page_height
                scale = min(scale_width, scale_height)

                matrix = pymupdf.Matrix(scale, scale)
                try:
                    pix = page.get_pixmap(matrix=matrix)
                except Exception as e:
                    raise RuntimeError(f'Failed to get the pixmap on page: {i}, error: {e}')

                png_path = outdir.joinpath(f'page_{i}.png')
                try:
                    pix.save(png_path)
                except Exception as e:
                    raise RuntimeError(f'Failed to save page: {i} as png, error: {e}')
                outlist_pngs.append(png_path)
            except Exception as e:
                raise Exception(f'Unexpected error while parsing pdf: {e}')
    finally:
        doc.close()

    if len(outlist_pngs) == 0:
        raise ValueError(f'Parsing pdf whent wrong we don\'t have any png from it...')

    return outlist_pngs


def _png_to_pptx(pngs_list: list[Path], pptx: Path) -> Path:
    print('Creating powerpoint')

    try:
        prs = Presentation()
    except Exception as e:
        raise Exception(f'Failed to init presention, error: {e}')

    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for img_path in pngs_list:
        try:
            if not img_path.exists():
                raise FileNotFoundError(f'png: "{img_path}" does not exists')

            blank_slide_layout = prs.slide_layouts[6]

            try:
                slide = prs.slides.add_slide(blank_slide_layout)
            except Exception as e:
                raise RuntimeError(f'Failed to add slide to presention, error: {e}')

            try:
                pic = slide.shapes.add_picture(str(img_path), 0, 0)
            except Exception as e:
                raise RuntimeError(f'Failed to add picture to slide, error: {e}')

            if pic.width <= 0 or pic.height <= 0:
                raise ValueError(f'png: {img_path}, have a invalid dimension')

            try:
                width_ratio = slide_width / pic.width
                height_ratio = slide_height / pic.height
                scaling_factor = min(width_ratio, height_ratio, 1)

                pic.width = int(pic.width * scaling_factor)
                pic.height = int(pic.height * scaling_factor)

                pic.left = int((slide_width - pic.width) / 2)
                pic.top = int((slide_height - pic.height) / 2)
            except Exception as e:
                raise RuntimeError(f'Failed to scale the image: {img_path}, error: {e}')
        except Exception as e:
            raise Exception(f'Unexpected error while creating powerpoint: {e}')

    if len(prs.slides) == 0:
        raise RuntimeError("Failed to create any slides in the presentation")

    try:
        prs.save(pptx.name)
    except Exception as e:
        raise IOError(f'Failed to save the powerpoint, error: {e}')

    return Path(pptx.name).resolve()


if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        prog='pdf2pptx',
        description='Convert pdf to pptx')
    parser.add_argument('filename', type=Path, help='The path to the pdf file')
    parser.add_argument('-r', '--resolution',
                        type=str,
                        default='1920x1080',
                        help='The target resolution of the presention (default 1920x1080)')
    args = parser.parse_args()

    args.filename = Path(str(args.filename).lower()).resolve()
    if not args.filename.exists():
        print(f'The pdf: "{args.filename}" does not exists', file=sys.stderr)
        exit(1)

    res = args.resolution.split('x')
    if len(res) != 2:
        print('Not engough arguments for resolution, use format: "<width>x<height>"',
              file=sys.stderr)
        exit(2)

    if not res[0].isdigit() or not res[1].isdigit():
        print(f'argument for resolution have a non digit in side it: {res}',
              file=sys.stderr)
        exit(3)

    RES = RESOLUTION(width=int(res[0]), height=int(res[1]))
    outname = args.filename.with_suffix('.pptx')

    main(pdf=args.filename.resolve(), reso=RES, out=outname)
